"""BIZ PATH 경진대회 PPT 생성 스크립트 — 24슬라이드."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BG      = RGBColor(0x0D, 0x1B, 0x2A)
ACCENT  = RGBColor(0x7B, 0xD0, 0x8D)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
MUTED   = RGBColor(0x8A, 0xA4, 0xB8)
YELLOW  = RGBColor(0xF5, 0xC5, 0x18)
RED     = RGBColor(0xE5, 0x53, 0x4B)
SURFACE = RGBColor(0x17, 0x27, 0x3A)
PURPLE  = RGBColor(0xC0, 0x8A, 0xF0)
BLUE    = RGBColor(0x7B, 0xB8, 0xD0)
DARK    = RGBColor(0x0A, 0x14, 0x1F)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(blank)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return s

def box(s, text, l, t, w, h, size=14, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True):
    b = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = b.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return b

def mbox(s, lines, l, t, w, h):
    b = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = b.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        txt, sz, bd, col = item
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(sz)
        r.font.bold = bd
        r.font.color.rgb = col

def rec(s, l, t, w, h, color):
    sh = s.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh

def divider(s, t, color=ACCENT):
    rec(s, 0.6, t, 12.1, 0.04, color)

def header(s, num, title, color=ACCENT):
    rec(s, 0, 0, 13.33, 0.72, SURFACE)
    box(s, num,   0.3,  0.1, 0.6, 0.52, size=18, bold=True, color=color)
    box(s, title, 0.95, 0.1, 11,  0.52, size=18, bold=True, color=WHITE)

def tbl(s, headers, rows, l, t, widths, rh=0.44):
    x = l
    for h, w in zip(headers, widths):
        rec(s, x, t, w - 0.03, rh, ACCENT)
        box(s, h, x + 0.06, t + 0.06, w - 0.14, rh - 0.1,
            size=10, bold=True, color=BG)
        x += w
    for ri, row in enumerate(rows):
        bg = SURFACE if ri % 2 == 0 else BG
        x = l
        rt = t + rh * (ri + 1)
        for ci, (cell, w) in enumerate(zip(row, widths)):
            rec(s, x, rt, w - 0.03, rh, bg)
            box(s, cell, x + 0.06, rt + 0.05, w - 0.14, rh - 0.08,
                size=10, color=ACCENT if ci == 0 else WHITE)
            x += w

def card(s, x, y, w, h, title, body, bar_color, title_size=13):
    rec(s, x, y, w, h, SURFACE)
    rec(s, x, y, w, 0.38, bar_color)
    box(s, title, x + 0.12, y + 0.06, w - 0.2, 0.3,
        size=title_size, bold=True, color=BG)
    box(s, body,  x + 0.12, y + 0.5,  w - 0.2, h - 0.6,
        size=11, color=WHITE, wrap=True)

def hint(s, text):
    box(s, f"💡  {text}", 0.6, 6.4, 12, 0.5, size=10, color=MUTED)


# ════════════════════════════════════════════════════════════
# 섹션 1. 제안 배경 및 출품작 소개
# ════════════════════════════════════════════════════════════

# ── 슬라이드 1: 표지 ─────────────────────────────────────────
s = slide()
rec(s, 0, 3.1, 13.33, 0.07, ACCENT)
box(s, "BIZ PATH", 1, 1.0, 11.33, 1.5,
    size=68, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
box(s, "창업 입지부터 상권 정책까지 — 흐름으로 읽는 서울 상권",
    1, 2.55, 11.33, 0.65, size=20, color=ACCENT, align=PP_ALIGN.CENTER)
box(s, '"왜 이 상권에서 창업해야 하는가. 데이터가 대답합니다."',
    1.5, 3.3, 10.33, 0.55, size=14, color=MUTED, align=PP_ALIGN.CENTER)
box(s, "2026 서울시 빅데이터 활용 경진대회  |  창업 부문",
    1, 6.6, 11.33, 0.5, size=12, color=MUTED, align=PP_ALIGN.CENTER)

# ── 슬라이드 2: 제안 배경 ────────────────────────────────────
s = slide()
header(s, "01", "제안 배경")
divider(s, 0.88)

# 인용구 박스
rec(s, 0.5, 1.02, 12.3, 1.1, DARK)
rec(s, 0.5, 1.02, 0.08, 1.1, ACCENT)
box(s, '"어디서 뭘 해야 망하지 않을까?"',
    0.75, 1.1, 11.5, 0.55, size=20, bold=True, color=WHITE)
box(s, "퇴직 후 창업을 준비하던 부모님이 반복하던 질문에서 BIZ PATH가 시작됐습니다.",
    0.75, 1.6, 11.5, 0.42, size=12, color=MUTED)

# 문제 3개 카드
problems = [
    (RED,    "정보 비대칭",
     "예비 창업자는 유동인구·폐업률·경쟁 분포를\n직접 수집·해석하기 어렵습니다.\n결국 임대 조건, 주변 조언, 개인 감각에 의존합니다."),
    (YELLOW, "중장년층 창업의 특수성",
     "퇴직 이후 창업은 단순한 도전이 아닙니다.\n실패는 노후 경제 안정과 직결됩니다.\n사전에 위험을 줄이는 객관적 근거가 필수입니다."),
    (BLUE,   "기존 도구의 한계",
     "공식 상권변화지표는 결과 스냅샷만 제공합니다.\n왜 나빠졌는지 원인 설명이 없고,\n흐름 기반 위험 신호는 포함되지 않습니다."),
]
for i, (col, title, body) in enumerate(problems):
    x = 0.5 + i * 4.2
    card(s, x, 2.25, 3.9, 3.85, title, body, col)

rec(s, 0.5, 6.2, 12.3, 0.95, SURFACE)
box(s,
    "BIZ PATH는 서울시 공공 빅데이터를 활용해 상권의 유입·유출 흐름, 폐업률, 네트워크 고립도를 분석하고,\n"
    "복잡한 상권 정보를 예비 창업자도 한 화면에서 직관적으로 확인할 수 있는 형태로 제공합니다.",
    0.72, 6.28, 11.7, 0.78, size=11, color=WHITE)

# ── 슬라이드 3: 출품작 소개 ──────────────────────────────────
s = slide()
header(s, "02", "출품작 소개")
box(s, "BIZ PATH — 공공데이터 기반 상권 위험 진단 플랫폼",
    0.6, 0.9, 12, 0.55, size=16, bold=True)
divider(s, 1.55)

features = [
    (ACCENT, "AI 창업 입지 추천",
     "업종·자치구 선택 → 추천·주의·비추천 9개 상권\n+ Claude AI 자연어 해설"),
    (YELLOW, "GRI 위험지수",
     "폐업률·순유출·고립도 결합 → 0~100 점수\n상권 침체 위험 사전 감지"),
    (BLUE,   "OD 이동 흐름 시각화",
     "실시간 인구 이동 흐름 애니메이션\n목적·시간대별 필터링"),
    (PURPLE, "3D 상권 시각화",
     "Deck.gl 3D 뷰 — GRI·순유입·폐업률 기준\n원기둥 높이·색상으로 위험도 표현"),
    (RED,    "흐름 단절 구간 시각화",
     "연결이 약해진 상권 구간 실 도로 경로 기반 표시\n'어디와의 연결이 끊겼는가' 직관 확인"),
    (MUTED,  "정책 추천 카드",
     "GRI 신호 → R4~R7 룰 자동 매칭\n규칙 기반 · 환각 없음 · CSV 출력"),
]
for i, (col, title, body) in enumerate(features):
    x = 0.5 + (i % 3) * 4.2
    y = 1.72 + (i // 3) * 2.42
    card(s, x, y, 3.9, 2.2, title, body, col, title_size=12)

# ── 슬라이드 4: 페르소나 ─────────────────────────────────────
s = slide()
header(s, "03", "페르소나")
box(s, "이런 분들을 위해 만들었습니다",
    0.6, 0.9, 12, 0.55, size=20, bold=True)
divider(s, 1.55)

rec(s, 0.5, 1.7, 5.9, 4.5, SURFACE)
rec(s, 0.5, 1.7, 5.9, 0.44, ACCENT)
box(s, "페르소나 A  |  예비 창업자", 0.72, 1.76, 5.4, 0.36,
    size=12, bold=True, color=BG)
mbox(s, [
    ("김미숙 · 53세 · 송파구 · 가정주부 25년 → 반찬가게 보조 2년", 12, True, WHITE),
    ("", 4, False, WHITE),
    ("강남구 어딘가에 반찬가게를 차릴까 망설이고 있다.", 11, False, MUTED),
    ("젊었을 때라면 실패해도 다시 시작하면 그만이었지만,", 11, False, MUTED),
    ("이제는 퇴직금이 전부다. 이 한 번이 마지막 기회다.", 11, True, WHITE),
    ("", 4, False, WHITE),
    ("• 인터넷 검색과 지인 조언만으로는 불안하다", 11, False, MUTED),
    ("• 어느 골목이 실제로 사람이 오가는지 모른다", 11, False, MUTED),
    ("• 비슷한 반찬가게가 얼마나 버티는지 알고 싶다", 11, False, MUTED),
    ("", 4, False, WHITE),
    ("필요한 것", 11, True, ACCENT),
    ('"강남 어느 상권에 들어가면 안 되는지\n숫자로 먼저 알려줘."', 11, False, WHITE),
], 0.72, 2.22, 5.4, 3.8)

rec(s, 6.9, 1.7, 5.9, 4.5, SURFACE)
rec(s, 6.9, 1.7, 5.9, 0.44, YELLOW)
box(s, "페르소나 B  |  소상공인 담당 공무원", 7.12, 1.76, 5.4, 0.36,
    size=12, bold=True, color=BG)
mbox(s, [
    ("박민준 · 41세 · 관악구 경제과 소속", 12, True, WHITE),
    ("", 4, False, WHITE),
    ("분기마다 관악구 전체 상권을 훑어 위험 상권을", 11, False, MUTED),
    ("직접 골라내야 하는데, 기준도 없고 시간도 없다.", 11, False, MUTED),
    ("", 4, False, WHITE),
    ("• GRI로 위험 상권 자동 식별 → 직접 뒤질 필요 없음", 11, False, WHITE),
    ("• 상권 클릭 → 폐업률·유동인구 즉시 확인", 11, False, WHITE),
    ("• 정책 카드(R4~R7) 자동 생성 → 대응 방향 즉시 파악", 11, False, WHITE),
    ("• CSV 출력 → 수기 작업 없이 결재 자료로 바로 전환", 11, False, WHITE),
    ("", 4, False, WHITE),
    ("필요한 것", 11, True, YELLOW),
    ('"어느 상권이 위험한지 자동으로 뽑아주고,\n이유까지 문서로 나오면 된다."', 11, False, WHITE),
], 7.12, 2.22, 5.4, 3.8)

# ── 슬라이드 5: 페르소나 C ───────────────────────────────────
s = slide()
header(s, "04", "페르소나 C  —  김미숙의 이야기")
divider(s, 0.88)

# 인용구
rec(s, 0.5, 1.02, 12.3, 0.95, DARK)
rec(s, 0.5, 1.02, 0.08, 0.95, ACCENT)
box(s, '"5천만원 까먹고 빈손으로 나오는 동네 아주머니가 되는 게 제일 무섭습니다."',
    0.75, 1.12, 11.5, 0.48, size=15, bold=True, color=WHITE)
box(s, "김미숙 · 53세 · 송파구 가락동 · 가정주부 25년 → 동네 반찬가게 보조 2년",
    0.75, 1.6, 11.5, 0.32, size=11, color=MUTED)

# 좌: 상황 카드
rec(s, 0.5, 2.1, 5.9, 4.55, SURFACE)
box(s, "상황", 0.72, 2.2, 5.4, 0.38, size=12, bold=True, color=ACCENT)
mbox(s, [
    ("강남구 어딘가에 반찬가게를 차릴까 망설이고 있다.", 11, False, WHITE),
    ("젊었을 때라면 실패해도 다시 시작하면 그만이었지만,", 11, False, WHITE),
    ("이제는 퇴직금이 전부다. 이 한 번이 마지막 기회다.", 11, True, WHITE),
    ("", 4, False, WHITE),
    ("• 인터넷 검색과 지인 조언만으로는 불안하다", 11, False, MUTED),
    ("• 어느 골목이 실제로 사람이 오가는지 모른다", 11, False, MUTED),
    ("• 비슷한 반찬가게가 얼마나 버티는지 알고 싶다", 11, False, MUTED),
    ("", 4, False, WHITE),
    ("강남이면 좋겠는데...", 11, True, YELLOW),
    ("• 매출은 크겠지만 임대료·권리금이 무섭다", 11, False, MUTED),
    ("관악이면 여유 있는데...", 11, True, BLUE),
    ("• 임대료는 괜찮지만 사람이 줄면 어떡하나", 11, False, MUTED),
    ("", 4, False, WHITE),
    ("딸이 말했다: \"엄마 이번엔 데이터 보고 정해\"", 11, True, WHITE),
], 0.72, 2.62, 5.4, 3.9)

# 우: BIZ PATH가 답해야 할 질문
rec(s, 6.9, 2.1, 5.9, 4.55, SURFACE)
box(s, "BIZ PATH가 답해야 할 질문", 7.12, 2.2, 5.4, 0.38,
    size=12, bold=True, color=YELLOW)
qna = [
    ('"여기 사람이 진짜 다니나?"',         "OD 흐름 · 유입 시계열"),
    ('"지금은 좋아도 곧 무너지는 거 아냐?"', "GRI 위험지수 · 흡수형/방출형 분류"),
    ('"임대료 더 세질 동네인가?"',          "강남 과열 신호 탐지"),
    ('"딸한테 어떻게 보여주지?"',           "CSV 다운로드"),
]
for i, (q, a) in enumerate(qna):
    y = 2.68 + i * 0.98
    rec(s, 7.1, y, 5.5, 0.85, BG)
    box(s, q, 7.25, y + 0.06, 5.2, 0.36, size=11, bold=True, color=WHITE)
    box(s, f"→ {a}", 7.25, y + 0.46, 5.2, 0.3, size=10, color=ACCENT)


# ════════════════════════════════════════════════════════════
# 섹션 2. 출품작 핵심내용
# ════════════════════════════════════════════════════════════

# ── 슬라이드 6(구 5): 공공데이터 활용 ───────────────────────
s = slide()
header(s, "04", "공공데이터 활용")
box(s, "서울시 공공데이터 5종을 결합했습니다",
    0.6, 0.9, 12, 0.55, size=18, bold=True)
divider(s, 1.55)

tbl(s,
    ["코드", "데이터셋", "출처", "활용"],
    [
        ["OA-14991", "서울 생활인구 (내국인)",          "서울 열린데이터광장", "시간대별 유동인구"],
        ["OA-15577", "우리마을가게 상권분석 (점포)",     "서울 열린데이터광장", "폐업률 · 점포 수"],
        ["OA-15572", "우리마을가게 상권분석 (추정매출)", "서울 열린데이터광장", "매출 기반 가중치"],
        ["OA-15576", "서울시 상권변화지표",              "서울 열린데이터광장", "B1 베이스라인 비교"],
        ["OA-22300", "수도권 광역 OD 기종점통행량",      "공공데이터포털",     "이동 흐름 핵심 데이터"],
    ],
    l=0.5, t=1.72, widths=[1.8, 4.0, 3.1, 3.4], rh=0.5)

rec(s, 0.5, 4.6, 12.3, 1.3, SURFACE)
box(s, "공간 결합 방식", 0.8, 4.68, 3, 0.4, size=12, bold=True, color=ACCENT)
box(s,
    "행정동 경계(행안부 SHP) × 상권 경계(서울시 SHP)  →  GeoPandas gpd.overlay() 공간 결합 (EPSG:5179 투영)\n"
    "→  adm_comm_mapping 테이블 생성  →  5개 데이터셋을 상권 단위로 통합하는 기반",
    0.8, 5.08, 11.5, 0.72, size=12, color=WHITE)

# ── 슬라이드 6: AI 창업 어드바이저 ──────────────────────────
s = slide()
header(s, "05", "핵심 기능 ①  AI 창업 입지 추천")
box(s, "추천만 주는 게 아닙니다. 왜 비추천인지도 알려드립니다.",
    0.6, 0.9, 12, 0.55, size=17, bold=True)
divider(s, 1.55)

mbox(s, [
    ("동작 방식", 13, True, ACCENT),
    ("", 4, False, WHITE),
    ("① 창업 희망 업종 + 관심 자치구 선택 후 '분석하기' 클릭", 12, False, WHITE),
    ("② GRI · 유동인구 · 폐업률 · 점포 수 · 중심성 5개 지표로 상권 점수화", 12, False, WHITE),
    ("③ 추천(상위 30%) / 주의(중간 40%) / 비추천(하위 30%) 분류", 12, False, WHITE),
    ("④ 각 등급 상위 3개씩, 총 9개 선별", 12, False, WHITE),
    ("⑤ Claude API가 수치 읽고 추천 이유·주의사항을 자연어로 생성", 12, False, WHITE),
], 0.6, 1.72, 5.6, 3.6)

tbl(s,
    ["지표", "가중치", "방향"],
    [
        ["GRI 위험도",        "0.25", "낮을수록 유리"],
        ["폐업률",            "0.25", "낮을수록 유리"],
        ["유동인구 (순유입)", "0.20", "높을수록 유리"],
        ["업종 점포 수",      "0.20", "중간값 선호 (경쟁 포화 반영)"],
        ["네트워크 중심성",   "0.10", "높을수록 유리"],
    ],
    l=6.5, t=1.72, widths=[2.9, 1.3, 2.5], rh=0.46)

hint(s, "화면 캡처: 추천·주의·비추천 카드 결과 화면  — 여기에 삽입")

# ── 슬라이드 7: GRI ──────────────────────────────────────────
s = slide()
header(s, "06", "핵심 기능 ②  GRI 상권 위험지수")
box(s, "매출이 멀쩡해도 흐름이 끊긴 상권을 미리 잡아냅니다",
    0.6, 0.9, 12, 0.55, size=17, bold=True)
divider(s, 1.55)

rec(s, 0.5, 1.72, 12.3, 1.0, SURFACE)
box(s, "GRI (Gentrification Risk Index)  —  0~100점, 높을수록 위험",
    0.8, 1.78, 10, 0.42, size=13, bold=True, color=ACCENT)
box(s, "폐업률 · 인구 순유출 · 네트워크 고립도 3가지를 z-score 정규화 후 가중 결합, percentile rank 0~100 변환",
    0.8, 2.13, 11.5, 0.5, size=12, color=WHITE)

rec(s, 0.5, 2.85, 12.3, 0.85, DARK)
box(s, "GRI  =  0.40 × z(폐업률)  +  0.33 × z(순유출)  +  0.27 × z(고립도)   →   percentile rank 0~100",
    0.8, 2.95, 11.5, 0.6, size=13, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)

mbox(s, [
    ("왜 필요한가?", 13, True, ACCENT),
    ("", 4, False, WHITE),
    ("• 매출만 보면 강남 압구정은 항상 '안전'으로 분류 — 흐름 신호를 놓침", 12, False, WHITE),
    ("• 유입 감소·연결성 약화를 매출 감소보다 먼저 포착해 사전 경고 가능", 12, False, WHITE),
    ("• H3 검증: GRI 상위 20% 상권의 다음 분기 폐업률이 하위 80%의 14.5배 (p=8×10⁻³¹)", 12, True, YELLOW),
], 0.6, 3.85, 12, 2.3)

hint(s, "화면 캡처: 상권 클릭 시 GRI 게이지 + 분기별 추이 차트  — 여기에 삽입")

# ── 슬라이드 8: OD 이동 흐름 ────────────────────────────────
s = slide()
header(s, "07", "핵심 기능 ③  인구 이동 흐름 기반 위험 탐지")
box(s, "단순 매출 지표로는 보이지 않는 흐름 기반 위험 신호를 감지합니다",
    0.6, 0.9, 12, 0.55, size=16, bold=True)
divider(s, 1.55)

mbox(s, [
    ("데이터 출처", 12, True, ACCENT),
    ("수도권 광역 OD 기종점통행량 (OA-22300) — 행정동 단위 → 상권 단위로 변환", 11, False, MUTED),
    ("", 5, False, WHITE),
    ("분석 방법 (NetworkX 유향 그래프)", 12, True, ACCENT),
    ("• 상권 간 이동 네트워크 구축 → net_flow, in/out_degree, centrality 산출", 12, False, WHITE),
    ("• 유입 감소 · 유출 증가 · 연결성 약화를 수치로 포착", 12, False, WHITE),
    ("", 5, False, WHITE),
    ("시각화", 12, True, ACCENT),
    ("• 출근·귀가·쇼핑·여가 목적별 이동 흐름 필터링", 12, False, WHITE),
    ("• 0~23시 시간대 슬라이더로 시간축 재생", 12, False, WHITE),
], 0.6, 1.72, 5.8, 4.5)

rec(s, 6.9, 1.72, 5.9, 4.5, SURFACE)
box(s, "창업자에게 주는 의미", 7.12, 1.84, 5.4, 0.42,
    size=13, bold=True, color=YELLOW)
mbox(s, [
    ('"오전 10시 이 상권에 쇼핑 목적 유입이 많다"', 12, True, WHITE),
    ("", 4, False, WHITE),
    ("→ 카페 운영 시간대 최적화 가능", 11, False, WHITE),
    ("→ 인근 직장인 점심 수요 파악 가능", 11, False, WHITE),
    ("→ 야간 영업 필요 여부 판단 가능", 11, False, WHITE),
    ("", 6, False, WHITE),
    ('"이 상권은 들어오는 사람보다 나가는 사람이 많다"', 12, True, WHITE),
    ("", 4, False, WHITE),
    ("→ 순유출 증가 = 상권 침체 조기 신호", 11, False, RED),
], 7.12, 2.35, 5.4, 3.7)

hint(s, "화면 캡처: OD 흐름 애니메이션 켜진 지도 화면  — 여기에 삽입")

# ── 슬라이드 9: 상권 상세 패널 ───────────────────────────────
s = slide()
header(s, "08", "핵심 기능 ④  상권 상세 패널")
box(s, "상권을 클릭하면 판단에 필요한 정보가 한 화면에 모입니다",
    0.6, 0.9, 12, 0.55, size=17, bold=True)
divider(s, 1.55)

detail_items = [
    (ACCENT, "GRI 위험 게이지",
     "0~100 점수 + 자치구 내 백분위\n분기별 추이 차트 (악화/개선 방향)"),
    (YELLOW, "유동인구 순유입",
     "OD 기반 해당 상권 순유입량\n시간대·목적별 분포"),
    (RED,    "전 폐업률",
     "상권 전 업종 기준 폐업률\n서울 평균 대비 위치 · 최근 4분기 추이"),
    (BLUE,   "업종 점포 현황",
     "동일 업종 점포 수 · 업종별 폐업률\n경쟁 밀도 파악"),
    (PURPLE, "창업 판단 요약  (AI 분석)",
     "총평 · 검토 업종 · 주의 업종\nClaude Haiku 기반 · 흐름+GRI 근거 포함"),
    (MUTED,  "지자체 대응 제안 (참고용)",
     "GRI 신호 기반 R4~R7 룰 자동 매칭\n수치 근거 포함 · 실제 시행 정책 아님"),
]
for i, (col, title, body) in enumerate(detail_items):
    x = 0.5 + (i % 2) * 6.2
    y = 1.72 + (i // 2) * 1.85
    card(s, x, y, 5.9, 1.65, title, body, col)

hint(s, "화면 캡처: 상권 클릭 시 상세 패널 전체 화면  — 여기에 삽입")

# ── 슬라이드 10: 흐름 단절 구간 시각화 ──────────────────────
s = slide()
header(s, "09", "핵심 기능 ⑤  흐름 단절 구간 시각화")
box(s, "'어느 상권이 위험한가'를 넘어 '어디와의 연결이 약해졌는가'를 보여줍니다",
    0.6, 0.9, 12, 0.55, size=15, bold=True)
divider(s, 1.55)

# 왼쪽: 3단계 위험도 + 단절 높으면 의미
box(s, "단절 위험 3단계", 0.6, 1.72, 6.0, 0.45, size=13, bold=True, color=ACCENT)
tbl(s,
    ["단계", "점수 기준", "색상", "상태"],
    [
        ["HIGH",   "0.75 이상",    "빨강 ●", "두 상권 연결이 거의 끊김 — 즉각 주목"],
        ["MEDIUM", "0.45 ~ 0.74", "분홍 ●", "흐름 약화 진행 중 — 추세 모니터링"],
        ["LOW",    "0.45 미만",   "보라 ●", "이상 징후 초기 — 경과 관찰"],
    ],
    0.6, 2.22, [1.1, 1.5, 1.1, 2.3])

mbox(s, [
    ("단절 위험이 높으면?", 13, True, RED),
    ("", 5, False, WHITE),
    ("두 상권 사이의 이동 흐름이 크게 약해진 상태입니다.", 11, False, WHITE),
    ("", 4, False, WHITE),
    ("→  해당 상권으로 유입되는 고객 공급이 감소 중", 11, True, YELLOW),
    ("→  GRI 위험도 상승 · 폐업률 증가로 이어질 수 있음", 11, False, WHITE),
    ("", 6, False, WHITE),
    ("선 굵기 = 영향받는 유동량에 비례 (최대 기준 1만 명)", 10, False, MUTED),
    ("굵을수록 더 많은 사람이 영향권에 있다는 의미입니다.", 10, False, MUTED),
], 0.6, 3.98, 6.0, 2.2)

# 오른쪽: 기존 차이 + 안정성
rec(s, 7.0, 1.72, 5.8, 4.5, SURFACE)
box(s, "기존 시각화와의 차이", 7.22, 1.84, 5.3, 0.42, size=13, bold=True, color=YELLOW)
mbox(s, [
    ("기존 방식", 11, True, MUTED),
    ("출발지 · 도착지를 직선으로 연결", 11, False, MUTED),
    ("지도 위에 떠 있는 선 → 실제 이동 경로와 무관", 11, False, MUTED),
    ("", 8, False, WHITE),
    ("BIZ PATH", 11, True, WHITE),
    ("ORS API로 실제 도로 경로를 계산", 11, False, WHITE),
    ("흐름 선이 실제 도로 위를 따라 표시됨", 11, True, ACCENT),
    ("", 10, False, WHITE),
    ("안정성 설계", 11, True, YELLOW),
    ("", 3, False, WHITE),
    ("ORS API 장애 시 자동 fallback → 캐시 재생", 11, False, WHITE),
    ("발표 중단 없이 정상 시연 가능", 11, False, WHITE),
], 7.22, 2.35, 5.3, 3.7)

hint(s, "화면 캡처: 흐름 단절 구간이 표시된 지도 화면  — 여기에 삽입")

# ── 슬라이드 11: 3D 상권 시각화 ──────────────────────────────
s = slide()
header(s, "10", "핵심 기능 ⑥  3D 상권 시각화")
box(s, "위험도를 높이와 색상으로 표현하는 입체 지도",
    0.6, 0.9, 12, 0.55, size=17, bold=True)
divider(s, 1.55)

mbox(s, [
    ("어떻게 동작하나", 13, True, ACCENT),
    ("", 4, False, WHITE),
    ("① 3D 모드 ON → 지도가 45도 기울어진 입체 시점으로 전환", 12, False, WHITE),
    ("② 각 상권 위치에 원기둥(CommerceColumnLayer)이 표시됨", 12, False, WHITE),
    ("③ 원기둥 높이와 색상으로 위험도를 직관적으로 파악", 12, False, WHITE),
], 0.6, 1.72, 6.0, 2.5)

tbl(s,
    ["기준 지표", "원기둥 의미"],
    [
        ["GRI 위험도",   "높을수록 위험한 상권 — 빨간색 + 높은 기둥"],
        ["순유입 인구",  "높을수록 유동인구가 많은 상권 — 초록색 + 높은 기둥"],
        ["폐업률",       "높을수록 폐업 위험 큰 상권 — 주황색 + 높은 기둥"],
        ["연결 중심성",  "높을수록 상권 네트워크에서 중심적 역할 — 파란색 + 높은 기둥"],
    ],
    l=0.6, t=4.1, widths=[3.5, 8.4], rh=0.5)

hint(s, "화면 캡처: 3D 상권 시각화 화면  — 여기에 삽입")

# ── 슬라이드 12: 기술성 ─────────────────────────────────────
s = slide()
header(s, "11", "기술성")
box(s, "생성형 AI · 네트워크 분석 · 통계를 결합했습니다",
    0.6, 0.9, 12, 0.55, size=18, bold=True)
divider(s, 1.55)

tech_items = [
    (ACCENT, "생성형 AI\nAnthropic Claude Haiku",
     "창업 어드바이저 자연어 해설 생성\n입력: 9개 상권 수치 데이터\n출력: 추천·비추천 이유 + 주의사항"),
    (YELLOW, "네트워크 분석\nNetworkX",
     "상권 간 OD 이동 유향 그래프 구축\nnet_flow · centrality 산출\nOD 데이터를 상권 단위로 면적 배분"),
    (BLUE,   "통계 분석\nscipy",
     "Pearson 상관 (H1 유입↔매출)\nWelch t-test (H3 GRI↔폐업)\nz-score 정규화 (GRI 합산)"),
]
for i, (col, title, body) in enumerate(tech_items):
    x = 0.5 + i * 4.2
    card(s, x, 1.72, 3.9, 4.0, title, body, col, title_size=12)


# ════════════════════════════════════════════════════════════
# 섹션 3. 기존 서비스와의 차별성
# ════════════════════════════════════════════════════════════

# ── 슬라이드 13: 차별점 ─────────────────────────────────────
s = slide()
header(s, "12", "기존 서비스와의 차별성")
box(s, "기존 지표가 '안전'이라 분류한 상권, BIZ PATH는 위험을 먼저 감지합니다",
    0.6, 0.9, 12, 0.55, size=16, bold=True)
divider(s, 1.55)

rec(s, 0.5, 1.72, 5.9, 4.9, SURFACE)
rec(s, 0.5, 1.72, 5.9, 0.44, ACCENT)
box(s, "비교 ①  서울시 공식 상권변화지표와 비교", 0.72, 1.78, 5.4, 0.36,
    size=11, bold=True, color=BG)
mbox(s, [
    ("공식 지표란?", 11, True, ACCENT),
    ("서울시가 분기마다 발표하는 상권 상태 등급.\n'정체·주의·상권쇠퇴·다이나믹' 4단계로 분류.", 10, False, MUTED),
    ("", 4, False, WHITE),
    ("비교 결과", 11, True, ACCENT),
    ("두 모델이 같은 상권을 위험으로 보는 비율: 58%", 10, False, WHITE),
    ("→ 공식 지표에는 없지만 BIZ PATH가 흐름 기반으로\n   추가 위험 신호를 감지한 상권:", 10, False, MUTED),
], 0.72, 2.24, 5.4, 2.3)
box(s, "187개", 1.2, 4.22, 3.5, 0.75, size=34, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)
box(s, "공식 지표 미포착이 아닌, 흐름 기준 추가 감지", 0.72, 5.05, 5.3, 0.38,
    size=10, color=ACCENT)

rec(s, 6.9, 1.72, 5.9, 4.9, SURFACE)
rec(s, 6.9, 1.72, 5.9, 0.44, BLUE)
box(s, "비교 ②  기존 매출 추세 모델과 비교", 7.12, 1.78, 5.4, 0.36,
    size=11, bold=True, color=BG)
mbox(s, [
    ("기존 매출 추세 모델이란?", 11, True, BLUE),
    ("직전 분기 매출 증감만으로 위험 상권을 판단하는 방식.\n매출이 멀쩡하면 '안전'으로 분류.", 10, False, MUTED),
    ("", 4, False, WHITE),
    ("비교 결과", 11, True, BLUE),
    ("두 모델이 같은 상권을 위험으로 보는 비율: 15%", 10, False, WHITE),
    ("→ 매출은 정상이지만 흐름이 끊긴 상권, BIZ PATH 추가 포착:", 10, False, MUTED),
], 7.12, 2.24, 5.4, 2.3)
box(s, "231개", 7.6, 4.22, 3.5, 0.75, size=34, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)
box(s, "매출이 아닌 '흐름'으로 보기 때문에 가능", 7.12, 5.05, 5.3, 0.38,
    size=10, color=BLUE)

# ── 슬라이드 14: 통계 검증 ──────────────────────────────────
s = slide()
header(s, "13", "통계 검증")
box(s, "1,650개 상권 데이터로 지표의 유효성을 직접 검증했습니다",
    0.6, 0.9, 12, 0.55, size=16, bold=True)
divider(s, 1.55)

# 행 1: H1 / H3
rec(s, 0.5, 1.68, 5.9, 1.85, SURFACE)
rec(s, 0.5, 1.68, 5.9, 0.38, BLUE)
box(s, "H1  |  유입이 많으면 매출도 높은가?", 0.72, 1.74, 5.4, 0.30,
    size=11, bold=True, color=BG)
mbox(s, [
    ("r = 0.106 / p = 2.83×10⁻⁵ / n = 1,565개 상권", 10, True, WHITE),
    ("", 3, False, WHITE),
    ("→ 약한 양의 상관 — 유입↑일 때 매출↑ 경향 확인", 10, False, WHITE),
    ("→ p = 0.00003 = 우연일 확률 0.003%  ✅ 통계 유의", 10, False, WHITE),
    ("단, r=0.11로 효과 크기 약함 — 매출엔 다른 요인도 작용", 9, False, MUTED),
], 0.72, 2.12, 5.4, 1.28)

rec(s, 6.9, 1.68, 5.9, 1.85, SURFACE)
rec(s, 6.9, 1.68, 5.9, 0.38, YELLOW)
box(s, "H3  |  GRI 위험 상권은 실제로 더 많이 폐업하는가?", 7.12, 1.74, 5.4, 0.30,
    size=11, bold=True, color=BG)
mbox(s, [
    ("14.5배 차이 / p = 8×10⁻³¹ / n = 1,650개 상권", 10, True, WHITE),
    ("", 3, False, WHITE),
    ("GRI 상위 20% vs 하위 80% — 폐업률 14.5배 차이", 10, False, WHITE),
    ("→ p = 8×10⁻³¹  ✅ GRI 예측력 강하게 입증", 10, True, YELLOW),
], 7.12, 2.12, 5.4, 1.28)

# 행 2: B1 / B3
rec(s, 0.5, 3.65, 5.9, 1.85, SURFACE)
rec(s, 0.5, 3.65, 5.9, 0.38, ACCENT)
box(s, "B1  |  서울시 공식 상권변화지표 대비", 0.72, 3.71, 5.4, 0.30,
    size=11, bold=True, color=BG)
mbox(s, [
    ("Jaccard = 0.157 / 추가 식별 187건 / 공통 탐지 143건", 10, True, WHITE),
    ("", 3, False, WHITE),
    ("공식 지표(OA-15576)가 잡지 못한 187개 상권을 추가 탐지", 10, False, WHITE),
    ("→ 흐름 기반 신호가 스냅샷 지표와 다른 위험을 포착", 10, True, ACCENT),
], 0.72, 4.09, 5.4, 1.28)

rec(s, 6.9, 3.65, 5.9, 1.85, SURFACE)
rec(s, 6.9, 3.65, 5.9, 0.38, PURPLE)
box(s, "B3  |  기존 매출 추세 모델 대비", 7.12, 3.71, 5.4, 0.30,
    size=11, bold=True, color=BG)
mbox(s, [
    ("Jaccard = 0.151 / 추가 위험 식별 231건", 10, True, WHITE),
    ("", 3, False, WHITE),
    ("매출 추세 모델과 거의 겹치지 않음 — 본질적으로 다른 신호", 10, False, WHITE),
    ("→ 강남 압구정·청담 등 '안정' 분류 상권도 조기 위험 탐지", 10, True, YELLOW),
], 7.12, 4.09, 5.4, 1.28)

# 요약
rec(s, 0.5, 5.6, 12.3, 0.82, DARK)
box(s, "검증 요약", 0.75, 5.67, 2.2, 0.35, size=11, bold=True, color=ACCENT)
box(s, "H3·B1·B3 모두 통과. H1 효과 약함·H2 표본 한계(n=39)는 투명 공개. 흐름 기반 탐지는 기존 지표와 본질적으로 다른 위험 신호를 포착합니다.",
    2.8, 5.67, 9.7, 0.68, size=10, color=WHITE)


# ════════════════════════════════════════════════════════════
# 섹션 4. 개발 과정 및 방법
# ════════════════════════════════════════════════════════════

# ── 슬라이드 15: 개발 타임라인 ───────────────────────────────
s = slide()
header(s, "14", "개발 과정 및 방법  —  개발 타임라인")
box(s, "5주 스프린트, 데이터 수집부터 AI 기능까지",
    0.6, 0.9, 12, 0.55, size=17, bold=True)
divider(s, 1.55)

weeks = [
    ("Week 1", "데이터 수집",
     "공공데이터 5종 수집 스크립트\nPostgreSQL+PostGIS 구축\n데이터 품질 검증"),
    ("Week 2", "분석 기반",
     "FastAPI 초기 구축\nNetworkX 그래프 분석\nGRI v1.0 공식 설계"),
    ("Week 3", "API 고도화",
     "API 5종 완성\nModule A/B/D/E 구현\n공간 결합 파이프라인"),
    ("Week 4", "시각화",
     "Deck.gl 흐름·3D 레이어\n실 도로 경로 단절 시각화\n통계 검증 5카드"),
    ("Week 5", "AI 완성",
     "Claude AI 창업 어드바이저\n배포 (Vercel + Railway)\n성능 최적화"),
]
colors = [MUTED, BLUE, ACCENT, YELLOW, PURPLE]
for i, (week, title, body) in enumerate(weeks):
    x = 0.4 + i * 2.52
    rec(s, x, 1.72, 2.42, 0.44, colors[i])
    box(s, week, x + 0.08, 1.78, 2.2, 0.32, size=11, bold=True, color=BG)
    rec(s, x + 1.15, 2.16, 0.12, 3.7, colors[i])
    rec(s, x, 2.24, 2.42, 3.55, SURFACE)
    box(s, title, x + 0.1, 2.32, 2.2, 0.42, size=12, bold=True, color=colors[i])
    box(s, body,  x + 0.1, 2.82, 2.2, 2.8,  size=10, color=WHITE)

# ── 슬라이드 16: 아키텍처 ────────────────────────────────────
s = slide()
header(s, "15", "개발 과정 및 방법  —  시스템 아키텍처")
box(s, "수집부터 시각화까지 완전 자동화 파이프라인",
    0.6, 0.9, 12, 0.55, size=17, bold=True)
divider(s, 1.55)

arch = [
    (MUTED,   "데이터 수집",  "공공데이터 API  (열린데이터광장 / 공공데이터포털)"),
    (ACCENT,  "저장 및 분석", "PostgreSQL + PostGIS + Redis  |  NetworkX · scipy · numpy"),
    (YELLOW,  "API 서버",     "FastAPI  —  /api/od/flows · /api/advisor · /api/barriers"),
    (BLUE,    "프론트엔드",   "React + Vite + MapLibre GL + Deck.gl (2D·3D 전환)"),
    (PURPLE,  "AI 해설",      "Anthropic Claude Haiku API  —  창업 어드바이저 자연어 생성"),
]
for i, (col, label, desc) in enumerate(arch):
    y = 1.72 + i * 0.94
    rec(s, 0.5, y, 9.6, 0.8, SURFACE)
    rec(s, 0.5, y, 2.1, 0.8, col)
    box(s, label, 0.6, y + 0.2, 1.9, 0.42, size=11, bold=True, color=BG)
    box(s, desc,  2.75, y + 0.2, 7.2, 0.42, size=11, color=WHITE)
    if i < len(arch) - 1:
        box(s, "↓", 5.25, y + 0.8, 0.4, 0.22, size=11, color=MUTED)

rec(s, 10.4, 1.72, 2.6, 4.7, SURFACE)
box(s, "배포", 10.6, 1.82, 2.2, 0.38, size=11, bold=True, color=ACCENT)
mbox(s, [
    ("프론트엔드", 9, True, MUTED), ("Vercel", 10, False, WHITE), ("", 4, False, WHITE),
    ("백엔드", 9, True, MUTED), ("Railway (Docker)", 10, False, WHITE), ("", 4, False, WHITE),
    ("데이터베이스", 9, True, MUTED), ("Supabase\nPostgreSQL+PostGIS", 10, False, WHITE),
], 10.6, 2.25, 2.3, 4.0)


# ════════════════════════════════════════════════════════════
# 섹션 5. IA
# ════════════════════════════════════════════════════════════

# ── 슬라이드 17: IA 화면 구조도 ─────────────────────────────
s = slide()
header(s, "16", "IA  —  화면 구조")
box(s, "3개 영역이 유기적으로 연결됩니다",
    0.6, 0.9, 12, 0.55, size=17, bold=True)
divider(s, 1.55)

rec(s, 0.4, 1.72, 3.5, 4.8, SURFACE)
rec(s, 0.4, 1.72, 3.5, 0.42, ACCENT)
box(s, "컨트롤 패널 (우측)", 0.58, 1.78, 3.2, 0.34, size=11, bold=True, color=BG)
mbox(s, [
    ("업종 선택", 10, True, ACCENT), ("창업 희망 업종 드롭다운", 10, False, WHITE), ("", 3, False, WHITE),
    ("자치구 필터", 10, True, ACCENT), ("복수 선택 · 전체/해제", 10, False, WHITE), ("", 3, False, WHITE),
    ("분기 선택", 10, True, ACCENT), ("2025Q1 ~ Q4", 10, False, WHITE), ("", 3, False, WHITE),
    ("AI 어드바이저 결과", 10, True, ACCENT), ("추천·주의·비추천 카드", 10, False, WHITE), ("", 3, False, WHITE),
    ("OD 흐름 제어", 10, True, ACCENT), ("목적·시간·강도 슬라이더", 10, False, WHITE),
], 0.58, 2.22, 3.2, 4.1)

rec(s, 4.1, 1.72, 5.1, 4.8, SURFACE)
rec(s, 4.1, 1.72, 5.1, 0.42, YELLOW)
box(s, "지도 (중앙)", 4.28, 1.78, 4.7, 0.34, size=11, bold=True, color=BG)
mbox(s, [
    ("MapLibre GL 기반 서울 지도", 10, False, MUTED), ("", 3, False, WHITE),
    ("상권 노드", 10, True, YELLOW), ("GRI 점수로 색상·크기 인코딩", 10, False, WHITE), ("", 3, False, WHITE),
    ("OD 흐름 레이어", 10, True, YELLOW), ("Deck.gl 파티클 애니메이션", 10, False, WHITE), ("", 3, False, WHITE),
    ("흐름 단절 레이어", 10, True, YELLOW), ("실 도로 경로 + 입자 효과", 10, False, WHITE), ("", 3, False, WHITE),
    ("3D 뷰 모드", 10, True, YELLOW), ("원기둥 높이·색상 위험도 표현", 10, False, WHITE),
], 4.28, 2.22, 4.8, 4.1)

rec(s, 9.4, 1.72, 3.5, 4.8, SURFACE)
rec(s, 9.4, 1.72, 3.5, 0.42, BLUE)
box(s, "상세 패널 (클릭 시)", 9.58, 1.78, 3.2, 0.34, size=11, bold=True, color=BG)
mbox(s, [
    ("GRI 게이지", 10, True, BLUE), ("0~100 + 백분위", 10, False, WHITE), ("", 3, False, WHITE),
    ("분기별 추이", 10, True, BLUE), ("TrendChart 시계열", 10, False, WHITE), ("", 3, False, WHITE),
    ("유동인구", 10, True, BLUE), ("순유입 · 목적별", 10, False, WHITE), ("", 3, False, WHITE),
    ("폐업률", 10, True, BLUE), ("업종별 비교", 10, False, WHITE), ("", 3, False, WHITE),
    ("지자체 대응 제안", 10, True, BLUE), ("R4~R7 자동 매칭", 10, False, WHITE),
], 9.58, 2.22, 3.2, 4.1)

for x in [3.6, 9.0]:
    box(s, "→", x, 3.9, 0.5, 0.5, size=20, bold=True, color=MUTED, align=PP_ALIGN.CENTER)

# ── 슬라이드 18: 사용자 흐름 ────────────────────────────────
s = slide()
header(s, "17", "IA  —  사용자 흐름 시나리오")
box(s, "예비 창업자 김미숙(53세, 송파구)이 강남 반찬가게 입지를 찾는 과정",
    0.6, 0.9, 12, 0.55, size=17, bold=True)
divider(s, 1.55)

steps = [
    ("①", "업종·자치구 선택",   "반찬가게 선택\n강남구 체크"),
    ("②", "AI 분석 실행",       "'분석하기' 클릭\n30초 내 결과"),
    ("③", "추천 카드 확인",     "추천 3개·주의 3개\n비추천 3개 + 이유"),
    ("④", "상권 클릭",          "지도에서 위치 확인\nGRI·폐업률 상세"),
    ("⑤", "OD 흐름 확인",      "오전 10시 쇼핑 유입\n시간대별 패턴"),
]
colors5 = [ACCENT, YELLOW, BLUE, PURPLE, MUTED]
for i, (num, title, body) in enumerate(steps):
    x = 0.4 + i * 2.52
    rec(s, x, 1.72, 2.42, 4.6, SURFACE)
    rec(s, x, 1.72, 2.42, 0.44, colors5[i])
    box(s, num,   x + 0.1,  1.78, 0.4, 0.34, size=14, bold=True, color=BG)
    box(s, title, x + 0.52, 1.78, 1.8, 0.34, size=11, bold=True, color=BG)
    box(s, body,  x + 0.1,  2.26, 2.2, 3.9,  size=11, color=WHITE)
    if i < 4:
        box(s, "→", x + 2.42, 3.7, 0.4, 0.44,
            size=18, bold=True, color=MUTED, align=PP_ALIGN.CENTER)

rec(s, 0.5, 6.45, 12.3, 0.72, SURFACE)
box(s, "결과: 강남 대치동 B상권 선택 → 유동인구 상위 20%, 반찬가게 폐업률 낮음, 귀가 유입 강함 → 입지 결정",
    0.72, 6.54, 11.8, 0.52, size=11, color=ACCENT)


# ════════════════════════════════════════════════════════════
# 섹션 6. 창업·사업화·매출·투자가능성
# ════════════════════════════════════════════════════════════

# ── 슬라이드 19: 시장 분석 ───────────────────────────────────
s = slide()
header(s, "18", "시장 분석")
box(s, "창업 상권 분석 수요는 크지만, 접근 가능한 도구가 없습니다",
    0.6, 0.9, 12, 0.55, size=16, bold=True)
divider(s, 1.55)

market = [
    (ACCENT, "서울시 연간\n신규 창업자", "약 10만 명\n(중소벤처기업부 2024)"),
    (YELLOW, "서울시\n소상공인 수",      "약 70만 명\n(통계청 2024)"),
    (BLUE,   "소상공인\n연간 폐업률",    "약 15%\n= 연 10만 명 폐업"),
    (RED,    "창업 전\n상권 분석 도구",  "전문 도구 부재\n공공데이터 수동 취합"),
]
for i, (col, title, val) in enumerate(market):
    x = 0.5 + (i % 2) * 6.2
    y = 1.72 + (i // 2) * 2.3
    rec(s, x, y, 5.9, 2.1, SURFACE)
    rec(s, x, y, 5.9, 0.42, col)
    box(s, title, x + 0.15, y + 0.07, 5.5, 0.34, size=12, bold=True, color=BG)
    box(s, val,   x + 0.15, y + 0.6,  5.5, 1.3,  size=20, bold=True, color=WHITE)

# ── 슬라이드 20: 수익 모델 ───────────────────────────────────
s = slide()
header(s, "19", "수익 모델")
box(s, "B2C + B2G 이중 수익 구조",
    0.6, 0.9, 12, 0.55, size=18, bold=True)
divider(s, 1.55)

rec(s, 0.5, 1.72, 5.9, 4.6, SURFACE)
rec(s, 0.5, 1.72, 5.9, 0.44, ACCENT)
box(s, "B2C  |  창업자 유료 리포트", 0.72, 1.78, 5.4, 0.36,
    size=13, bold=True, color=BG)
mbox(s, [
    ("대상", 11, True, ACCENT), ("예비 창업자, 소상공인", 11, False, WHITE), ("", 4, False, WHITE),
    ("상품", 11, True, ACCENT), ("상권 상세 분석 리포트 (PDF)\nGRI + 경쟁 현황 + AI 해설 포함", 11, False, WHITE), ("", 4, False, WHITE),
    ("가격", 11, True, ACCENT), ("건당 2~5만원", 11, False, WHITE), ("", 4, False, WHITE),
    ("추산 매출", 11, True, ACCENT), ("연 신규 창업자 10만 × 구매율 5%\n= 5,000건 × 3만원  →  연 1.5억원", 11, False, WHITE),
], 0.72, 2.28, 5.4, 3.8)

rec(s, 6.9, 1.72, 5.9, 4.6, SURFACE)
rec(s, 6.9, 1.72, 5.9, 0.44, YELLOW)
box(s, "B2G  |  기관 구독", 7.12, 1.78, 5.4, 0.36,
    size=13, bold=True, color=BG)
mbox(s, [
    ("대상", 11, True, YELLOW), ("소상공인진흥공단\n각 자치구 경제과\n창업지원센터", 11, False, WHITE), ("", 4, False, WHITE),
    ("상품", 11, True, YELLOW), ("API 연동 + 대시보드 구독\n위험 상권 자동 리포트 + CSV", 11, False, WHITE), ("", 4, False, WHITE),
    ("가격", 11, True, YELLOW), ("자치구 연 500만원 구독", 11, False, WHITE), ("", 4, False, WHITE),
    ("추산 매출", 11, True, YELLOW), ("25개 자치구 × 500만원\n= 연 1.25억원\n(공단 계약 시 규모 확대)", 11, False, WHITE),
], 7.12, 2.28, 5.4, 3.8)

# ── 슬라이드 21: 확장 가능성 + 투자 포인트 ──────────────────
s = slide()
header(s, "20", "확장 가능성 및 투자 포인트")
divider(s, 0.9)

phases = [
    (ACCENT, "단기  v1.1",
     "서울 전역 25개 자치구 확장\n(동일 파이프라인, 데이터 재수집)\nGRI 4항목 복원 (임대료 추가)"),
    (YELLOW, "중기  v2.0",
     "전국 광역시 확장\n정책 담당자 전용 대시보드\n모바일 현장 조회"),
    (PURPLE, "장기  비전",
     "정부 디지털플랫폼정부 정책 연계\n소상공인 지원 데이터 인프라로 통합\n\"창업 전 필수 확인 플랫폼\""),
]
for i, (col, phase, body) in enumerate(phases):
    x = 0.4 + i * 4.2
    card(s, x, 1.05, 3.9, 3.2, phase, body, col)

rec(s, 0.5, 4.45, 12.3, 2.55, SURFACE)
box(s, "투자 포인트", 0.8, 4.55, 3, 0.42, size=13, bold=True, color=YELLOW)
mbox(s, [
    ("• 공공데이터 기반  →  데이터 조달 비용 사실상 0", 12, False, WHITE),
    ("• AI·자동화로 분석 운영 인건비 최소화", 12, False, WHITE),
    ("• B2G 시장  →  안정적 장기 계약 기반 (자치구·공단)", 12, False, WHITE),
    ("• 서울 강남·관악 MVP 검증 완료  →  전국 확장 리스크 낮음", 12, False, WHITE),
], 0.8, 5.02, 11.5, 1.8)

# ── 슬라이드 22: 활용방안 ────────────────────────────────────
s = slide()
header(s, "21", "활용방안")
box(s, "예비 창업자, 행정기관, 연구기관이 함께 쓰는 의사결정 도구",
    0.6, 0.9, 12, 0.55, size=15, bold=True)
divider(s, 1.55)

users = [
    (ACCENT, "예비 창업자",
     "• AI 어드바이저로 업종별 후보 상권 비교\n"
     "• 점포 계약 전 GRI·폐업률·유동인구 확인\n"
     "• CSV 출력 → 가족·상담사와 함께 검토\n"
     "• 임대 조건만 보고 결정하는 위험 감소"),
    (YELLOW, "자치구청 / 소상공인 지원기관",
     "• GRI 기반 위험 상권 선제 모니터링\n"
     "• 상권 활성화 예산 배분 근거 자료\n"
     "• 골목상권 지원사업 후보 선정\n"
     "• 컨설팅·금융지원 대상 발굴\n"
     "• /api/export/csv → 결재 문서 즉시 전환"),
    (BLUE, "도시재생 전문가 / 연구기관",
     "• 흐름 단절 구간으로 보행 환경 개선 후보지 검토\n"
     "• /api/barriers, /api/barrier-routes 활용\n"
     "• 상권 간 접근성 저하 구간 분석\n"
     "• 재생 우선 후보지 선정 근거 자료"),
]
for i, (col, title, body) in enumerate(users):
    x = 0.4 + i * 4.2
    card(s, x, 1.72, 3.95, 5.0, title, body, col)

# ── 슬라이드 23: 기대효과 ────────────────────────────────────
s = slide()
header(s, "22", "기대효과")
box(s, "폐업 후 대응에서 위험 신호 사전 발견으로 상권 지원 체계를 전환합니다",
    0.6, 0.9, 12, 0.55, size=14, bold=True)
divider(s, 1.55)

effects = [
    (ACCENT, "창업 실패 비용 감소",
     "GRI 위험 상권을 사전에 피해 창업하면\n폐업 확률을 낮출 수 있습니다.\n\n"
     "H3 검증: 위험 상권 폐업률 = 안전 상권의 14.5배\n\n"
     "입지 선택 단계에서의 잘못된 결정이\n가장 큰 창업 실패 원인임을 데이터로 확인했습니다."),
    (YELLOW, "행정 사전 개입 가능",
     "GRI는 다음 분기 폐업을 미리 가리키는\n정책 선별 지표로 활용 가능합니다.\n\n"
     "폐업 발생 후 대응 → 위험 신호 선발견 후 개입\n\n"
     "규칙 기반 정책 카드(R4~R7)와 CSV 출력으로\n수기 보고 시간 단축, 근거 자료 즉시 생성."),
    (BLUE, "분석 투명성 확보",
     "5개 검증 항목(H1·H2·H3, B1·B3)을\n화면에서 직접 공개합니다.\n\n"
     "한계(H1 효과 약함, H2 표본 소규모)도 함께 명시해\n맹목적 신뢰가 아닌 근거 기반 판단을 지원합니다.\n\n"
     "정책 카드는 생성형 AI 없이 규칙 기반으로 생성 →\n환각 없음, 감사·보고 과정에서 추적 가능."),
]
for i, (col, title, body) in enumerate(effects):
    x = 0.4 + i * 4.2
    card(s, x, 1.72, 3.95, 5.0, title, body, col)


# ════════════════════════════════════════════════════════════
# 섹션 7. 개발 툴 및 참고문헌
# ════════════════════════════════════════════════════════════

# ── 슬라이드 24: 개발 툴 및 참고문헌 ────────────────────────
s = slide()
header(s, "23", "개발 툴 및 참고문헌")
divider(s, 0.9)

rec(s, 0.5, 1.05, 7.8, 3.6, SURFACE)
box(s, "개발 툴 및 기술 스택", 0.72, 1.14, 7, 0.4, size=12, bold=True, color=ACCENT)
tbl(s,
    ["분류", "기술"],
    [
        ["백엔드",      "Python · FastAPI · SQLAlchemy · PostgreSQL+PostGIS · Redis"],
        ["프론트엔드",  "React · Vite · TypeScript · MapLibre GL · Deck.gl"],
        ["분석",        "NetworkX · scipy · numpy · pandas · geopandas"],
        ["AI",          "Anthropic Claude Haiku (claude-haiku-4-5-20251001)"],
        ["배포",        "Vercel (프론트) · Railway Docker (백엔드) · Supabase (DB)"],
    ],
    l=0.5, t=1.58, widths=[2.0, 5.78], rh=0.42)

rec(s, 8.5, 1.05, 4.5, 3.6, SURFACE)
box(s, "공공데이터 출처", 8.7, 1.14, 4.1, 0.4, size=12, bold=True, color=YELLOW)
mbox(s, [
    ("OA-14991  서울 생활인구", 10, False, WHITE),
    ("OA-15577  상권분석 (점포)", 10, False, WHITE),
    ("OA-15572  상권분석 (매출)", 10, False, WHITE),
    ("OA-15576  상권변화지표", 10, False, WHITE),
    ("OA-22300  수도권 광역 OD", 10, False, WHITE),
    ("", 3, False, WHITE),
    ("행정동 경계 SHP  (행정안전부)", 10, False, MUTED),
    ("상권 경계 SHP  (서울 열린데이터광장)", 10, False, MUTED),
], 8.7, 1.62, 4.1, 2.9)

rec(s, 0.5, 4.82, 12.3, 2.1, SURFACE)
box(s, "참고문헌", 0.72, 4.9, 3, 0.38, size=12, bold=True, color=BLUE)
mbox(s, [
    ("• 서울 열린데이터광장  https://data.seoul.go.kr", 10, False, WHITE),
    ("• 공공데이터포털  https://www.data.go.kr", 10, False, WHITE),
    ("• 중소벤처기업부 소상공인 통계 2024", 10, False, WHITE),
    ("• Anthropic Claude API Documentation  https://docs.anthropic.com", 10, False, WHITE),
    ("• NetworkX Documentation  https://networkx.org", 10, False, WHITE),
], 0.72, 5.35, 11.5, 1.45)


out = "docs/spiceMap_경진대회_PPT.pptx"
prs.save(out)
print(f"저장 완료: {out}  ({len(prs.slides)}슬라이드)")
