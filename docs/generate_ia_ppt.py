"""Generate the spiceMap IA PowerPoint deck.

This script uses python-pptx so the output is a normal PowerPoint file rather
than a hand-written OOXML package.
"""
from __future__ import annotations

import sys
from pathlib import Path


try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt
except ModuleNotFoundError:
    temp_targets = [Path(r"C:\tmp\codex_pptx_user"), Path(r"C:\tmp\codex_pptx")]
    for temp_target in temp_targets:
        if temp_target.exists():
            sys.path.insert(0, str(temp_target))
            break
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches, Pt
    except ModuleNotFoundError:
        raise SystemExit(
            "python-pptx is required. Install it with: "
            "python -m pip install python-pptx"
        )


OUT = Path("docs/spicemap_IA_screen_structure.pptx")

BG = RGBColor(0x0D, 0x1B, 0x2A)
SURFACE = RGBColor(0x17, 0x27, 0x3A)
DARK = RGBColor(0x0A, 0x14, 0x1F)
ACCENT = RGBColor(0x7B, 0xD0, 0x8D)
BLUE = RGBColor(0x60, 0xA5, 0xFA)
YELLOW = RGBColor(0xFB, 0xBF, 0x24)
RED = RGBColor(0xF8, 0x71, 0x71)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0xA6, 0xB4, 0xC2)
BORDER = RGBColor(0x24, 0x38, 0x4D)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


def set_bg(slide, color=BG) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: int = 14,
    color=WHITE,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    fill=None,
    border=None,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border is not None:
        shape.line.color.rgb = border
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()

    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.08)
    frame.margin_top = Inches(0.04)
    frame.margin_bottom = Inches(0.04)

    for i, line in enumerate(text.split("\n")):
        para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        para.alignment = align
        run = para.add_run()
        run.text = line
        run.font.name = "Malgun Gothic"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return shape


def add_rect(slide, x: float, y: float, w: float, h: float, fill, border=None, rounded=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = border
        shape.line.width = Pt(1)
    return shape


def add_card(slide, x: float, y: float, w: float, h: float, title: str, body: str, bar_color):
    add_rect(slide, x, y, w, h, SURFACE, BORDER)
    add_rect(slide, x, y, w, 0.42, bar_color, rounded=True)
    add_text(slide, title, x + 0.12, y + 0.06, w - 0.24, 0.3, 12, BG, True)
    add_text(slide, body, x + 0.15, y + 0.55, w - 0.3, h - 0.65, 11, WHITE)


def add_header(slide, num: str, title: str):
    add_rect(slide, 0, 0, 13.333, 0.72, SURFACE, rounded=False)
    add_text(slide, num, 0.35, 0.12, 0.65, 0.42, 18, ACCENT, True)
    add_text(slide, title, 1.0, 0.12, 11.8, 0.42, 18, WHITE, True)
    add_rect(slide, 0.6, 0.88, 12.1, 0.04, ACCENT, rounded=False)


def add_slide(title: str):
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    return slide


def build() -> None:
    slide = add_slide("IA Overview")
    add_text(slide, "spiceMap IA 화면 구조", 0.8, 0.85, 11.7, 0.8, 40, WHITE, True, PP_ALIGN.CENTER)
    add_text(slide, "필터설정 → AI 분석 → 상권 탐색", 1.3, 1.65, 10.7, 0.55, 24, ACCENT, True, PP_ALIGN.CENTER)
    steps = [
        ("1. 필터설정", "업종 선택\n관심 자치구 선택\n분석 범위 확정", ACCENT),
        ("2. AI 분석", "5개 지표 가중합산\n3등급 자동분류\nClaude AI 자연어 해설", BLUE),
        ("3. 상권 탐색", "추천 카드 클릭\n지도 위치 확인\n시간대·유형별 유동인구 확인", YELLOW),
    ]
    for i, (title, body, color) in enumerate(steps):
        x = 0.75 + i * 4.15
        add_card(slide, x, 2.65, 3.55, 2.7, title, body, color)
        if i < 2:
            add_text(slide, "→", x + 3.58, 3.55, 0.5, 0.4, 30, MUTED, True, PP_ALIGN.CENTER)
    add_text(
        slide,
        "창업자가 조건을 입력하면 AI가 상권을 선별하고, 지도와 유동인구로 최종 판단을 돕는 단일 화면 구조",
        1.1,
        6.25,
        11.2,
        0.55,
        14,
        MUTED,
        align=PP_ALIGN.CENTER,
    )

    slide = add_slide("Main Screen Layout")
    add_header(slide, "01", "Main Screen Layout")
    add_rect(slide, 0.65, 1.25, 7.25, 5.75, DARK, BORDER)
    add_text(slide, "지도 영역", 0.95, 1.55, 6.7, 0.35, 18, WHITE, True)
    add_text(
        slide,
        "• 서울 상권 지도\n• 추천/주의/비추천 등급 마커\n• 선택 상권 위치 강조\n• 시간대·유형별 유동인구 흐름 레이어",
        0.95,
        2.1,
        6.55,
        1.5,
        14,
        MUTED,
    )
    for x, y, label, color in [(2.1, 4.35, "추천", BLUE), (3.95, 3.75, "주의", YELLOW), (5.8, 4.75, "비추천", RED)]:
        add_rect(slide, x, y, 1.6, 0.7, color)
        add_text(slide, label, x + 0.05, y + 0.13, 1.5, 0.35, 14, BG, True, PP_ALIGN.CENTER)
    add_rect(slide, 8.25, 1.25, 4.45, 5.75, SURFACE, BORDER)
    add_text(slide, "분석 패널", 8.55, 1.55, 3.9, 0.35, 18, WHITE, True)
    add_card(slide, 8.55, 2.1, 3.85, 1.15, "1. 필터설정", "업종, 관심 자치구, 분석 실행", ACCENT)
    add_card(slide, 8.55, 3.45, 3.85, 1.35, "2. AI 분석", "점수 계산, 등급 분류, Claude 해설, 추천 카드", BLUE)
    add_card(slide, 8.55, 5.0, 3.85, 1.35, "3. 상권 탐색", "카드 클릭, 지도 선택, 상세 지표, 유동인구", YELLOW)

    slide = add_slide("Step 1 Filter")
    add_header(slide, "02", "Step 1. 필터설정")
    add_text(slide, "사용자가 분석 범위를 먼저 좁히는 단계", 0.7, 1.1, 12, 0.45, 20, WHITE, True)
    add_card(slide, 0.7, 1.8, 3.8, 3.9, "업종 선택", "• 창업 업종 드롭다운\n• /api/advisor/industries 기반\n• 선택값: industry_nm\n• 미선택 시 분석 버튼 비활성", ACCENT)
    add_card(slide, 4.75, 1.8, 3.8, 3.9, "관심 자치구 선택", "• 서울 25개 구 선택\n• 전체 선택 / 초기화\n• 권역 그룹 / 검색\n• 선택값: districts", BLUE)
    add_card(slide, 8.8, 1.8, 3.8, 3.9, "분석 실행", "• 선택 조건으로 POST 요청\n• quarter는 보조 필터\n• 로딩/에러 상태 표시\n• 결과 패널로 이어짐", YELLOW)
    add_text(slide, "핵심 IA 원칙: 첫 화면에서 사용자가 해야 할 일은 '업종과 지역을 고르는 것'으로 단순화한다.", 0.95, 6.35, 11.6, 0.45, 13, MUTED)

    slide = add_slide("Step 2 AI Analysis")
    add_header(slide, "03", "Step 2. AI 분석")
    add_text(slide, "5개 지표를 하나의 창업 적합도 점수로 합산하고, 3등급으로 자동분류", 0.7, 1.06, 12, 0.45, 18, WHITE, True)
    rows = [
        ("GRI 역점수", "25%", "상권 위험도가 낮을수록 유리"),
        ("유동인구량", "20%", "선택 조건의 유입 규모"),
        ("폐업 안정성", "25%", "폐업률이 낮을수록 유리"),
        ("업종 점포 기반", "20%", "해당 업종 시장 기반"),
        ("네트워크 중심성", "10%", "연결성과 접근성"),
    ]
    for i, (metric, weight, meaning) in enumerate(rows):
        y = 1.85 + i * 0.72
        add_rect(slide, 0.75, y, 7.1, 0.55, SURFACE if i % 2 == 0 else DARK, BORDER, rounded=False)
        add_text(slide, metric, 0.9, y + 0.1, 2.2, 0.3, 12, WHITE, True)
        add_text(slide, weight, 3.3, y + 0.1, 0.8, 0.3, 12, ACCENT, True, PP_ALIGN.CENTER)
        add_text(slide, meaning, 4.3, y + 0.1, 3.25, 0.3, 11, MUTED)
    add_card(slide, 8.25, 1.85, 4.2, 1.25, "추천", "점수 상위 30%\n우선 검토 상권", BLUE)
    add_card(slide, 8.25, 3.35, 4.2, 1.25, "주의", "중간 40%\n조건부 검토 상권", YELLOW)
    add_card(slide, 8.25, 4.85, 4.2, 1.25, "비추천", "하위 30%\n리스크 우선 확인", RED)
    add_text(slide, "Claude AI는 점수 결과를 요약, 주의 문구, 카드별 자연어 이유로 변환한다.", 0.8, 6.4, 11.8, 0.45, 13, MUTED)

    slide = add_slide("Step 3 Explore")
    add_header(slide, "04", "Step 3. 상권 탐색")
    add_text(slide, "추천 카드에서 지도 검토로 이어지는 탐색 단계", 0.7, 1.06, 12, 0.45, 18, WHITE, True)
    add_card(slide, 0.7, 1.75, 3.3, 3.85, "추천 카드 클릭", "• 상권명 / 자치구\n• 추천 점수 / 등급\n• 핵심 이유 / 리스크\n• 다음 행동 제안", BLUE)
    add_text(slide, "→", 4.05, 3.25, 0.5, 0.4, 30, MUTED, True, PP_ALIGN.CENTER)
    add_card(slide, 4.55, 1.75, 3.3, 3.85, "지도 위치 확인", "• 선택 상권 강조\n• 주변 상권 비교\n• 등급 마커 유지\n• 지도 중심 이동", ACCENT)
    add_text(slide, "→", 7.9, 3.25, 0.5, 0.4, 30, MUTED, True, PP_ALIGN.CENTER)
    add_card(slide, 8.4, 1.75, 3.9, 3.85, "유동인구 확인", "• 0~23시 시간대 슬라이더\n• 유형별 OD 흐름 필터\n• 흐름 밀도 조절\n• 재생/일시정지", YELLOW)
    add_text(slide, "탐색 결과: 사용자는 AI 추천을 맹신하지 않고, 위치·지표·시간대별 유동 흐름으로 최종 판단한다.", 1.0, 6.35, 11.4, 0.45, 13, MUTED, align=PP_ALIGN.CENTER)

    slide = add_slide("State Data Flow")
    add_header(slide, "05", "IA State & Data Flow")
    add_text(slide, "입력값, 분석 결과, 지도 상태의 연결", 0.7, 1.08, 12, 0.4, 18, WHITE, True)
    blocks = [
        ("입력 상태", "industry_nm\nselectedDistricts\nselectedQuarter", ACCENT),
        ("AI 분석 API", "/api/advisor/startup\n점수 계산\n등급 분류\nClaude 해설", BLUE),
        ("결과 상태", "summary / caution\n추천 카드 목록\ntierMap", YELLOW),
        ("지도 상태", "selectedNode\n등급 마커\n유동인구 레이어", RED),
    ]
    for i, (title, body, color) in enumerate(blocks):
        x = 0.6 + i * 3.15
        add_card(slide, x, 2.0, 2.75, 3.0, title, body, color)
        if i < len(blocks) - 1:
            add_text(slide, "→", x + 2.78, 3.15, 0.35, 0.35, 24, MUTED, True, PP_ALIGN.CENTER)
    add_rect(slide, 0.75, 5.85, 11.85, 0.85, SURFACE, BORDER)
    add_text(slide, "PPT 메시지: spiceMap IA는 '조건 입력 → AI 판단 → 지도 검증'의 한 방향 흐름으로 설계되어 발표자가 화면 시연 순서를 설명하기 쉽다.", 0.95, 6.05, 11.45, 0.35, 12, WHITE)


if __name__ == "__main__":
    build()
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"created {OUT} ({len(prs.slides)} slides)")
