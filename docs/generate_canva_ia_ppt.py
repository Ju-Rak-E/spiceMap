"""Generate a Canva-import-friendly PPTX for the spiceMap IA.

Canva often flattens SVG uploads into a single image. PPTX imports usually keep
text boxes, rectangles, and lines editable, so this deck is built from native
PowerPoint shapes rather than embedded images.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


OUT = Path("docs/spicemap_ia_canva_editable.pptx")

WIDE_W = 13.333
WIDE_H = 7.5

TEXT = RGBColor(0x11, 0x18, 0x27)
MUTED = RGBColor(0x6B, 0x72, 0x80)
LINE = RGBColor(0x8B, 0x98, 0xA8)
ROOT = RGBColor(0xF9, 0x73, 0x16)
ROOT_LINE = RGBColor(0xC2, 0x41, 0x0C)
SECTION = RGBColor(0xFB, 0xBF, 0x24)
SECTION_LINE = RGBColor(0xD9, 0x77, 0x06)
NODE = RGBColor(0xFF, 0xFF, 0xFF)
LEAF = RGBColor(0xF8, 0xFA, 0xFC)
BOX_LINE = RGBColor(0xA7, 0xB0, 0xBB)
LEAF_LINE = RGBColor(0xCB, 0xD5, 0xE1)


prs = Presentation()
prs.slide_width = Inches(WIDE_W)
prs.slide_height = Inches(WIDE_H)
blank = prs.slide_layouts[6]


def add_line(slide, x1, y1, x2, y2):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    line.line.color.rgb = LINE
    line.line.width = Pt(1.1)
    return line


def add_textbox(slide, text, x, y, w, h, size=10, bold=False, color=TEXT, align=PP_ALIGN.CENTER):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.background()
    box.line.fill.background()
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    frame.margin_left = Inches(0.03)
    frame.margin_right = Inches(0.03)
    frame.margin_top = Inches(0.01)
    frame.margin_bottom = Inches(0.01)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Malgun Gothic"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_box(slide, cx, y, w, h, label, kind="node"):
    if kind == "root":
        fill, stroke, size, bold, color = ROOT, ROOT_LINE, 12, True, RGBColor(0xFF, 0xFF, 0xFF)
    elif kind == "section":
        fill, stroke, size, bold, color = SECTION, SECTION_LINE, 12, True, TEXT
    elif kind == "leaf":
        fill, stroke, size, bold, color = LEAF, LEAF_LINE, 8.5, False, TEXT
    else:
        fill, stroke, size, bold, color = NODE, BOX_LINE, 9.5, True, TEXT

    x = cx - w / 2
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = stroke
    shape.line.width = Pt(1)

    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    frame.margin_left = Inches(0.03)
    frame.margin_right = Inches(0.03)
    frame.margin_top = Inches(0.01)
    frame.margin_bottom = Inches(0.01)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = label
    run.font.name = "Malgun Gothic"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def add_header(slide, screen_no, subtitle):
    add_textbox(slide, f"spiceMap IA 화면 구조 {screen_no}", 0, 0.12, WIDE_W, 0.3, 19, True)
    add_textbox(slide, subtitle, 0, 0.42, WIDE_W, 0.22, 9, False, MUTED)


def add_root(slide, left_label, right_label):
    root_cx = WIDE_W / 2
    add_box(slide, root_cx, 0.82, 1.55, 0.36, "spiceMap", "root")
    add_line(slide, root_cx, 1.18, root_cx, 1.48)
    add_box(slide, root_cx, 1.48, 1.55, 0.34, "메인 화면", "node")

    branch_y = 2.18
    section_y = 2.45
    left_cx = 3.35
    right_cx = 9.98
    add_line(slide, root_cx, 1.82, root_cx, branch_y)
    add_line(slide, left_cx, branch_y, right_cx, branch_y)

    for cx, label in [(left_cx, left_label), (right_cx, right_label)]:
        add_line(slide, cx, branch_y, cx, section_y)
        add_box(slide, cx, section_y, 1.95, 0.38, label, "section")
    return left_cx, right_cx, section_y


def add_group(slide, cx, xoff, y, title, leaves):
    group_cx = cx + xoff
    node_w, node_h = 1.78, 0.31
    leaf_w, leaf_h = 1.58, 0.23
    leaf_gap = 0.08
    add_box(slide, group_cx, y, node_w, node_h, title, "node")
    first_leaf_y = y + node_h + 0.16
    add_line(slide, group_cx, y + node_h, group_cx, first_leaf_y)
    for idx, leaf in enumerate(leaves):
        ly = first_leaf_y + idx * (leaf_h + leaf_gap)
        if idx > 0:
            add_line(slide, group_cx, ly - leaf_gap, group_cx, ly)
        add_box(slide, group_cx, ly, leaf_w, leaf_h, leaf, "leaf")


def add_section_groups(slide, section_cx, section_y, groups):
    positions = [
        (-1.25, 3.2),
        (1.25, 3.2),
        (-1.25, 5.25),
        (1.25, 5.25),
    ]
    for (title, leaves), (xoff, y) in zip(groups, positions):
        group_cx = section_cx + xoff
        joint_y = y - 0.18
        add_line(slide, section_cx, section_y + 0.38, section_cx, joint_y)
        add_line(slide, section_cx, joint_y, group_cx, joint_y)
        add_line(slide, group_cx, joint_y, group_cx, y)
        add_group(slide, section_cx, xoff, y, title, leaves)


def add_footer(slide, text):
    add_textbox(slide, text, 0.8, 7.12, 11.75, 0.18, 8.8, False, MUTED)


def build_slide(screen_no, subtitle, left_label, right_label, left_groups, right_groups, footer):
    slide = prs.slides.add_slide(blank)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    add_header(slide, screen_no, subtitle)
    left_cx, right_cx, section_y = add_root(slide, left_label, right_label)
    add_section_groups(slide, left_cx, section_y, left_groups)
    add_section_groups(slide, right_cx, section_y, right_groups)
    add_footer(slide, footer)


screen1_left = [
    ("지도 기본 레이어", ["서울 행정구역", "상권 경계", "기본 지도 스타일"]),
    ("AI 결과 표시", ["추천/주의/비추천 마커", "선택 상권 강조", "상권 등급 색상"]),
    ("유동인구 레이어", ["OD 흐름선", "시간대별 흐름", "유형별 흐름"]),
    ("지도 인터랙션", ["줌 / 이동", "상권 클릭 선택", "지도 중심 이동"]),
]
screen1_right = [
    ("업종 선택", ["업종 목록 조회", "업종 드롭다운", "선택 업종 표시"]),
    ("관심 자치구 선택", ["전체 선택 / 초기화", "권역 그룹 선택", "자치구 검색", "구별 선택 버튼"]),
    ("분석 기준", ["분기 선택", "선택 조건 요약", "분석 가능 상태"]),
    ("분석 실행", ["AI 분석 버튼", "로딩 상태", "오류 메시지"]),
]
screen2_left = [
    ("분석 데이터", ["상권 분석 데이터", "업종 점포 데이터", "OD 유동인구"]),
    ("5개 지표 가중합산", ["GRI 역점수 25%", "유동인구량 20%", "폐업 안정성 25%", "업종 점포 기반 20%", "네트워크 중심성 10%"]),
    ("3등급 자동분류", ["추천: 상위 30%", "주의: 중간 40%", "비추천: 하위 30%"]),
    ("Claude AI 해설", ["전체 요약", "주의 문구", "카드별 추천 이유"]),
]
screen2_right = [
    ("추천 카드 목록", ["상권명 / 자치구", "점수 / 등급", "핵심 이유 / 리스크", "다음 행동 제안"]),
    ("지도 위치 확인", ["선택 상권 위치", "지도 중심 이동", "주변 상권 비교"]),
    ("상권 상세 지표", ["GRI", "순유동", "폐업률", "적합도 / 상권 특성"]),
    ("유동인구 확인", ["시간대 슬라이더", "유형별 필터", "흐름 밀도 조절", "재생 / 일시정지"]),
]


if __name__ == "__main__":
    build_slide(
        "1/2",
        "지도 영역 + 필터설정",
        "지도 영역",
        "1. 필터설정",
        screen1_left,
        screen1_right,
        "화면 1: 사용자가 분석 범위를 정하고 지도에서 기본 상권 구조를 확인하는 영역",
    )
    build_slide(
        "2/2",
        "AI 분석 + 상권 탐색",
        "2. AI 분석",
        "3. 상권 탐색",
        screen2_left,
        screen2_right,
        "화면 2: AI가 추천 결과를 만들고 사용자가 추천 상권을 지도·지표·유동인구로 탐색하는 영역",
    )
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"created {OUT} ({len(prs.slides)} slides)")

